"""Extract chart data from specific slides as proper JSON.

Slide 3:  New students by education level (LINE_MARKERS, 32-year time series)
Slide 9:  Academic staff by degree level (COLUMN_STACKED)
Slide 19: Patent applications (COLUMN_CLUSTERED)
"""
import json
import sys
from pathlib import Path
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

PPTX = "KMUTT Long Term Trend 2568-(@19052569).pptx"
OUT_DIR = Path("web/public/data")
OUT_DIR.mkdir(parents=True, exist_ok=True)

prs = Presentation(PPTX)

def chart_to_dict(chart, slide_idx):
    """Generic chart -> dict. Captures series name, categories, values."""
    data = {
        "slide": slide_idx,
        "chart_type": str(chart.chart_type),
        "has_title": chart.has_title,
        "title": chart.chart_title.text_frame.text if chart.has_title else None,
        "series": [],
    }
    try:
        cats = []
        for plot in chart.plots:
            if plot.categories:
                cats = [str(c) for c in plot.categories]
                break
        data["categories"] = cats
    except Exception as e:
        data["categories"] = []
        data["cat_error"] = str(e)

    for plot in chart.plots:
        for series in plot.series:
            data["series"].append({
                "name": series.name if series.name else "",
                "values": [None if v is None else float(v) for v in series.values],
            })
    return data

# Slides are 0-indexed in python-pptx
TARGETS = {3: "students-new", 9: "faculty-degree", 19: "patents"}

for slide_idx, fname in TARGETS.items():
    slide = prs.slides[slide_idx - 1]
    chart_shape = None
    for shape in slide.shapes:
        if shape.has_chart:
            chart_shape = shape
            break
    if not chart_shape:
        print(f"!! Slide {slide_idx}: no chart found")
        continue
    raw = chart_to_dict(chart_shape.chart, slide_idx)
    out_path = OUT_DIR / f"{fname}.raw.json"
    out_path.write_text(json.dumps(raw, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Slide {slide_idx} -> {out_path}")
    print(f"  type: {raw['chart_type']}")
    print(f"  categories ({len(raw['categories'])}): {raw['categories'][:10]}{'...' if len(raw['categories'])>10 else ''}")
    print(f"  series ({len(raw['series'])}):")
    for s in raw['series']:
        non_null = [v for v in s['values'] if v is not None]
        print(f"    - {s['name']!r}  len={len(s['values'])}  non-null={len(non_null)}  first 5: {s['values'][:5]}")
    print()
