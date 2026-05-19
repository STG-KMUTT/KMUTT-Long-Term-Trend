"""Build bilingual chart JSON files from the PPTX.

Numbers come straight from PPTX. Translations and metadata are kept here so
they survive PPT re-extraction. Run this whenever PPT updates.
"""
import json
import sys
from pathlib import Path
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

PPTX = "KMUTT Long Term Trend 2568-(@19052569).pptx"
OUT_DIR = Path("web/src/data")
OUT_DIR.mkdir(parents=True, exist_ok=True)

# Slide -> bilingual annotations + series styling
SPECS = {
    3: {
        "id": "students-new",
        "section": "education",
        "chart_type": "line",
        "title": {"th": "จำนวนนักศึกษาใหม่", "en": "New Students Enrolled"},
        "subtitle": {
            "th": "จำแนกตามระดับการศึกษา",
            "en": "By education level",
        },
        "series_meta": {
            "ปริญญาตรี":     {"key": "bachelor", "name_en": "Bachelor's",     "color": "#f29400"},
            "บัณฑิตศึกษา":   {"key": "graduate", "name_en": "Graduate (Master's + Ph.D.)", "color": "#1e6091"},
            "รวม":           {"key": "total",    "name_en": "Total",          "color": "#0f172a", "emphasis": True},
        },
        "methodology": {
            "th": "ข้อมูลปี 2546-2561 เป็นข้อมูล ณ ภาคการศึกษาที่ 1 (ภาคที่รับนักศึกษา) ตั้งแต่ปี 2562 เป็นข้อมูล ณ ภาคการศึกษาที่ 2 รับจำนวนนักศึกษาตามสถานภาพ ปี 2568 มีนักศึกษาระดับอนุปริญญา (ป.) เพิ่มอีก 78 คน",
            "en": "2003-2018 data is from semester 1 (admission semester). From 2019 onwards, data is from semester 2 based on enrollment status. 2025 also includes 78 students at the Diploma level.",
        },
        "source": {
            "th": "หนังสือสารสนเทศ ปีการศึกษา 2536-2568, กลุ่มงานวิจัยสถาบันและสารสนเทศ สำนักงานยุทธศาสตร์ มจธ.",
            "en": "Information Books 1993-2025, Institutional Research & Information Unit, Strategy Office, KMUTT",
        },
    },
    9: {
        "id": "faculty-degree",
        "section": "personnel",
        "chart_type": "stacked-bar",
        "title": {"th": "จำนวนบุคลากรสายวิชาการ จำแนกตามวุฒิการศึกษา",
                  "en": "Academic Staff by Highest Degree"},
        "subtitle": {"th": "อาจารย์และนักวิจัยของ มจธ.",
                     "en": "Faculty and researchers at KMUTT"},
        "series_meta": {
            "ปริญญาตรี":   {"key": "bachelor", "name_en": "Bachelor's", "color": "#fcd34d"},
            "ปริญญาโท":    {"key": "master",   "name_en": "Master's",   "color": "#f59e0b"},
            "ปริญญาเอก":   {"key": "doctorate","name_en": "Doctorate",  "color": "#b45309"},
            "รวม":         {"key": "total",    "name_en": "Total",      "color": "#0f172a", "exclude_from_stack": True},
        },
        "methodology": {
            "th": "ตั้งแต่ปี 2547 วุฒิการศึกษาของอาจารย์รวมสถานภาพพนักงานมหาวิทยาลัยลักษณะพิเศษ (พนักงานมหาวิทยาลัยลูกจ้างมหาวิทยาลัยอาจารย์ชาวต่างชาติ) และตั้งแต่ปี 2549 มีการนับรวมพนักงานพิเศษวิจัยที่เป็นบุคลากรสายวิชาการ",
            "en": "From 2004, academic staff degree data includes university staff (university lecturers, foreign faculty). From 2006, special research staff classified as academic staff are also included.",
        },
        "source": {
            "th": "หนังสือสารสนเทศในปีการศึกษานั้น ๆ, กลุ่มงานวิจัยสถาบันและสารสนเทศ สำนักงานยุทธศาสตร์ ข้อมูลปฐมภูมิจากสำนักงานบริหารทรัพยากรบุคคล",
            "en": "Information Books for the respective academic year, Institutional Research & Information Unit, Strategy Office. Primary data from HR Office.",
        },
    },
    19: {
        "id": "patents",
        "section": "research",
        "chart_type": "clustered-bar",
        "title": {"th": "การยื่นขอจดสิทธิบัตรและอนุสิทธิบัตร",
                  "en": "Patent and Petty Patent Filings"},
        "subtitle": {"th": "ผลงานทรัพย์สินทางปัญญาของ มจธ.",
                     "en": "Intellectual property output at KMUTT"},
        "series_meta": {
            "สิทธิบัตรที่ยื่นขอ":       {"key": "patent_filed",      "name_en": "Patents filed",         "color": "#f29400"},
            "สิทธิบัตรที่ได้รับ":       {"key": "patent_granted",    "name_en": "Patents granted",       "color": "#d97f00"},
            "อนุสิทธิบัตรที่ยื่นขอ":     {"key": "petty_filed",       "name_en": "Petty patents filed",   "color": "#1e6091"},
            "อนุสิทธิบัตรที่ได้รับ":     {"key": "petty_granted",     "name_en": "Petty patents granted", "color": "#0c4a6e"},
            "ยื่นขอสะสม":              {"key": "filed_cumulative",  "name_en": "Cumulative filed",      "color": "#94a3b8", "is_cumulative": True},
            "ได้รับสะสม":              {"key": "granted_cumulative","name_en": "Cumulative granted",    "color": "#475569", "is_cumulative": True},
        },
        "methodology": {
            "th": "รับข้อมูลจากหน่วยงานที่เกี่ยวข้อง มีการปรับตัวเลขย้อนหลังเนื่องจากมีการยกเลิกหรือเปลี่ยนคำขอสิทธิบัตรหรืออนุสิทธิบัตรในบางปี ซึ่งส่งผลให้จำนวนการยื่นขอเปลี่ยนแปลงจากรายงานครั้งก่อน",
            "en": "Data received from relevant offices. Numbers may be retroactively adjusted due to cancelled or modified patent/petty-patent applications, which can change filing counts compared to prior reports.",
        },
        "source": {
            "th": "หนังสือรายงานประจำปี มจธ., กลุ่มงานวิจัยสถาบันและสารสนเทศ สำนักงานยุทธศาสตร์ ข้อมูลปฐมภูมิจากสำนักงานวิจัยนวัตกรรมและพันธมิตร",
            "en": "KMUTT Annual Reports, Institutional Research & Information Unit, Strategy Office. Primary data from the Research, Innovation & Partnerships Office.",
        },
    },
}

prs = Presentation(PPTX)

for slide_idx, spec in SPECS.items():
    slide = prs.slides[slide_idx - 1]
    chart_shape = next((s for s in slide.shapes if s.has_chart), None)
    if not chart_shape:
        print(f"!! Slide {slide_idx}: no chart")
        continue
    chart = chart_shape.chart
    cats = []
    for plot in chart.plots:
        if plot.categories:
            cats = [str(c) for c in plot.categories]
            break

    out_series = []
    for plot in chart.plots:
        for series in plot.series:
            raw_name = series.name or ""
            meta = spec["series_meta"].get(raw_name)
            if meta is None:
                print(f"  WARN slide {slide_idx}: unknown series {raw_name!r}, skipping")
                continue
            values = [None if v is None else float(v) for v in series.values]
            entry = {
                "key": meta["key"],
                "name": {"th": raw_name, "en": meta["name_en"]},
                "color": meta["color"],
                "values": values,
            }
            for opt in ("emphasis", "exclude_from_stack", "is_cumulative"):
                if meta.get(opt):
                    entry[opt] = True
            out_series.append(entry)

    out = {
        "id": spec["id"],
        "slide": slide_idx,
        "section": spec["section"],
        "chart_type": spec["chart_type"],
        "title": spec["title"],
        "subtitle": spec["subtitle"],
        "categories_buddhist": cats,
        "series": out_series,
        "methodology": spec["methodology"],
        "source": spec["source"],
    }
    out_path = OUT_DIR / f"{spec['id']}.json"
    out_path.write_text(json.dumps(out, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Wrote {out_path}  ({len(cats)} cats × {len(out_series)} series)")
