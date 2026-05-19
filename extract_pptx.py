from pptx import Presentation
import sys

# Force stdout to UTF-8
sys.stdout.reconfigure(encoding='utf-8')

pptx_file = "KMUTT Long Term Trend 2568-(@19052569).pptx"
prs = Presentation(pptx_file)

def fix_thai(text):
    """Some PPTX runs come back as mojibake when fonts use legacy encoding.
    Try to recover by re-encoding latin-1 -> cp874 (Thai)."""
    try:
        return text.encode('latin-1', errors='strict').decode('cp874', errors='strict')
    except (UnicodeEncodeError, UnicodeDecodeError):
        return text

print(f"Total slides: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides, 1):
    print(f"\n========== Slide {i} ==========")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs)
                if text.strip():
                    print(fix_thai(text))
        if shape.has_table:
            print("  [TABLE]")
            for row in shape.table.rows:
                cells = [fix_thai(cell.text.strip()) for cell in row.cells]
                print("    " + " | ".join(cells))
        if shape.shape_type == 13:
            print(f"  [IMAGE: {shape.name}]")
        if shape.has_chart:
            ch = shape.chart
            print(f"  [CHART: type={ch.chart_type}]")
            try:
                for plot in ch.plots:
                    for series in plot.series:
                        vals = list(series.values)
                        cats = list(plot.categories) if plot.categories else []
                        name = fix_thai(series.name) if series.name else ""
                        print(f"    Series: {name}")
                        for c, v in zip(cats, vals):
                            print(f"      {fix_thai(str(c))}: {v}")
            except Exception as e:
                print(f"    (chart read err: {e})")
